import os
import sys
import json
import getpass
import pandas as pd
from pptx import Presentation
from moviepy import VideoFileClip
import whisper

# Windowsローカル環境用のFFmpegパス指定を復活
os.environ["PATH"] += os.pathsep + r"C:\Users\s_ise\Desktop\調査用\knowledge-base-rag\ffmpeg-2026-06-15-git-44d082edc8-essentials_build\bin"

# ==========================================
# 1. 自動設定
# ==========================================
username = getpass.getuser()

if len(sys.argv) > 1:
    TARGET_DIR = os.path.abspath(sys.argv[1])
else:
    TARGET_DIR = f"C:\\Users\\{username}\\Desktop\\ナレッジ変換"

TARGET_FOLDER_NAME = os.path.basename(TARGET_DIR)
OUTPUT_BASE_DIR = os.path.join(TARGET_DIR, TARGET_FOLDER_NAME)

OFFICE_EXTENSIONS = ('.xlsx', '.xlsm', '.pptx', '.ppt')
VIDEO_EXTENSIONS = ('.mp4', '.m4a', '.mp3', '.wav')

# ==========================================
# 2. 変換・文字起こし用関数
# ==========================================
def convert_excel_to_json(file_path, output_path):
    """Excelの全シートを読み込み、構造化JSONを出力"""
    try:
        excel_data = pd.read_excel(file_path, sheet_name=None)
        result = {}
        for sheet_name, df in excel_data.items():
            df_filled = df.fillna("")  # 空白セルを空文字に置換
            result[sheet_name] = df_filled.to_dict(orient="records")
            
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(result, f, ensure_ascii=False, indent=2)
        print("   -> [成功] ExcelをJSONに変換しました。")
        return True
    except Exception as e:
        print(f"   -> [失敗] Excel変換エラー: {e}")
        return False

def convert_pptx_to_json(file_path, output_path):
    """パワポのスライドごとにテキストを構造化してJSONを出力"""
    try:
        prs = Presentation(file_path)
        slides_data = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            slides_data.append({
                "slide_number": i + 1,
                "content": "\n".join(slide_text)
            })
            
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(slides_data, f, ensure_ascii=False, indent=2)
        print("   -> [成功] パワポをJSONに変換しました。")
        return True
    except Exception as e:
        print(f"   -> [失敗] パワポ変換エラー: {e}")
        return False

def transcribe_video(file_path, model, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    base_name, ext = os.path.splitext(os.path.basename(file_path))
    audio_path = os.path.join(output_dir, base_name + "_temp.mp3")
    txt_path = os.path.join(output_dir, base_name + ".txt")
    
    try:
        if ext.lower() == '.mp4':
            video = VideoFileClip(file_path)
            if video.audio is None:
                print("   -> [スキップ] 音声が含まれていない動画です。")
                return False
            video.audio.write_audiofile(audio_path, logger=None)
            target_audio = audio_path
        else:
            target_audio = file_path

        result = model.transcribe(target_audio, language="ja")
        with open(txt_path, "w", encoding="utf-8") as f:
            f.write(result["text"])
        
        print(f"   -> [成功] テキストファイルを生成しました: {os.path.basename(txt_path)}")
        if os.path.exists(audio_path):
            os.remove(audio_path)
        return True
    except Exception as e:
        print(f"   -> [失敗] 文字起こしエラー: {e}")
        if os.path.exists(audio_path):
            os.remove(audio_path)
        return False

# ==========================================
# 3. メイン処理ロジック
# ==========================================
def clean_folder(target_path):
    if not os.path.exists(target_path):
        print(f"【エラー】指定されたフォルダが見つかりません: {target_path}")
        return

    print("-> 文字起こしAI (Whisper) を初期化中...")
    model = whisper.load_model("base")

    print(f"\n--- 変換＆文字起こし処理を開始します ---")
    print(f"対象フォルダ: {target_path}")
    print(f"出力先フォルダ: {OUTPUT_BASE_DIR}\n")

    for root, dirs, files in os.walk(target_path):
        dirs[:] = [d for d in dirs if d != TARGET_FOLDER_NAME and not d.startswith('.')]

        relative_path = os.path.relpath(root, target_path)
        current_output_dir = os.path.abspath(os.path.join(OUTPUT_BASE_DIR, relative_path))

        for file in files:
            file_path = os.path.join(root, file)
            file_lower = file.lower()

            if file.startswith('~$') or file.startswith('.~lock.'):
                continue

            # スクリプト自身を誤って処理しないよう対策
            if file == "clean.py":
                continue

            base_name, _ = os.path.splitext(file)
            expected_json = os.path.join(current_output_dir, base_name + ".json")

            # ----------------------------------------------------
            # ① Excel / パワポを見つけたらJSON変換
            # ----------------------------------------------------
            if file_lower.endswith(OFFICE_EXTENSIONS):
                if os.path.exists(expected_json):
                    print(f"   -> [スキップ] すでにJSONが存在します: {file}")
                    continue

                os.makedirs(current_output_dir, exist_ok=True)
                print(f"【JSON変換】Officeファイルを解析します: {file}")
                
                if file_lower.endswith(('.xlsx', '.xlsm')):
                    convert_excel_to_json(file_path, expected_json)
                elif file_lower.endswith(('.pptx', '.ppt')):
                    convert_pptx_to_json(file_path, expected_json)
                continue

            # ----------------------------------------------------
            # ② 動画 / 音声を見つけたら文字起こし
            # ----------------------------------------------------
            if file_lower.endswith(VIDEO_EXTENSIONS):
                expected_txt = os.path.join(current_output_dir, base_name + ".txt")
                if os.path.exists(expected_txt):
                    print(f"   -> [スキップ] すでにテキストが存在します: {file}")
                    continue

                print(f"【文字起こし】動画/音声をテキスト化します: {file}")
                transcribe_video(file_path, model, current_output_dir)
                continue

if __name__ == "__main__":
    clean_folder(TARGET_DIR)
    print("\n--- すべての処理が完了しました！ ---")